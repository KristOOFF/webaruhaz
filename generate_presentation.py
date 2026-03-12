"""
NeoCoffee Webáruház - Prezentáció generáló
Generál egy teljes PPTX bemutatót a projektről
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Színpaletta ──────────────────────────────────────────────────────────────
COFFEE_DARK   = RGBColor(0x1A, 0x0A, 0x02)   # nagyon sötét barna
COFFEE_BROWN  = RGBColor(0x4E, 0x27, 0x0D)   # mély kávébarna
COFFEE_MID    = RGBColor(0x7B, 0x44, 0x18)   # közepes barna
COFFEE_LIGHT  = RGBColor(0xC4, 0x8B, 0x4A)   # latte arany
CREAM         = RGBColor(0xF5, 0xE6, 0xCB)   # krémszín
PINK_ACCENT   = RGBColor(0xEC, 0x48, 0x99)   # pink akcentus
PURPLE_ACCENT = RGBColor(0x93, 0x33, 0xEA)   # lila akcentus
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xE5, 0xE7, 0xEB)
DARK_BG       = RGBColor(0x0D, 0x0E, 0x1A)   # sötét kék-fekete háttér
SLATE_MID     = RGBColor(0x47, 0x55, 0x69)
GREEN_OK      = RGBColor(0x22, 0xC5, 0x5E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # üres elrendezés


# ─────────────────────────────────────────────────────────────────────────────
# SEGÉDFÜGGVÉNYEK
# ─────────────────────────────────────────────────────────────────────────────

def add_bg(slide, fill_color=DARK_BG, gradient=True):
    """Teljes diahátter téglalap (sötét kék-fekete, opcionálisan gradiens)."""
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = fill_color
    bg.line.fill.background()
    return bg


def add_rect(slide, x, y, w, h, color, alpha_workaround=False):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=20, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_multiline_text(slide, lines, x, y, w, h, font_size=14,
                       bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                       line_colors=None):
    """Több sorból álló szövegdoboz, soronként eltérő formázás."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if line_colors and i < len(line_colors):
            run.font.color.rgb = line_colors[i]
        else:
            run.font.color.rgb = color
    return txb


def gradient_bar(slide, y=0.55, h=0.08):
    """Vékony pink-lila vízszintes gradient vonal."""
    # two-tone workaround: két szomszédos téglalap
    bar1 = slide.shapes.add_shape(1,
        0, Inches(y), prs.slide_width // 2, Inches(h))
    bar1.fill.solid(); bar1.fill.fore_color.rgb = PINK_ACCENT
    bar1.line.fill.background()

    bar2 = slide.shapes.add_shape(1,
        prs.slide_width // 2, Inches(y), prs.slide_width // 2, Inches(h))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = PURPLE_ACCENT
    bar2.line.fill.background()


def add_pill(slide, text, x, y, w=1.6, h=0.35,
             bg=COFFEE_MID, fg=WHITE, font_size=11):
    """Lekerekített pill badge."""
    pill = slide.shapes.add_shape(
        9,  # rounded rectangle
        Inches(x), Inches(y), Inches(w), Inches(h))
    pill.fill.solid()
    pill.fill.fore_color.rgb = bg
    pill.line.color.rgb = bg
    tf = pill.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = fg
    return pill


def section_header(slide, title, subtitle="", bar_y=0.55):
    add_bg(slide, DARK_BG)
    gradient_bar(slide, bar_y)
    add_text(slide, title,
             0.4, 0.7, 12.5, 1.1,
             font_size=36, bold=True, color=CREAM, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 0.4, 1.65, 12.5, 0.5,
                 font_size=16, bold=False, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – BORÍTÓ
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, DARK_BG)

# Dekor körök (glassmorphism hangulathoz)
c1 = slide.shapes.add_shape(9, Inches(-1.2), Inches(-1.2), Inches(4), Inches(4))
c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x2D, 0x1A, 0x08)
c1.line.fill.background()

c2 = slide.shapes.add_shape(9, Inches(10.5), Inches(4.5), Inches(3.5), Inches(3.5))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x1E, 0x12, 0x30)
c2.line.fill.background()

gradient_bar(slide, y=0.0, h=0.06)

# Glassmorphism kártya a cím mögé
card = slide.shapes.add_shape(9,
    Inches(1.2), Inches(1.5), Inches(10.9), Inches(4.5))
card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x22, 0x14, 0x0A)
card.line.color.rgb = RGBColor(0x7B, 0x44, 0x18)

# Cím
add_text(slide, "☕  NeoCoffee",
         1.5, 1.7, 10.5, 1.4,
         font_size=52, bold=True, color=COFFEE_LIGHT,
         align=PP_ALIGN.CENTER)

add_text(slide, "Webáruház",
         1.5, 2.85, 10.5, 0.8,
         font_size=38, bold=True, color=CREAM,
         align=PP_ALIGN.CENTER)

add_text(slide, "Modern kávé webshop – SvelteKit · Express.js · SQLite",
         1.5, 3.65, 10.5, 0.55,
         font_size=17, bold=False, color=LIGHT_GRAY,
         align=PP_ALIGN.CENTER)

# Tech badge-ek
badges = [
    ("SvelteKit", COFFEE_MID),
    ("TypeScript", RGBColor(0x31, 0x78, 0xC6)),
    ("Tailwind CSS", RGBColor(0x06, 0x8F, 0x8A)),
    ("Express.js", RGBColor(0x26, 0x32, 0x38)),
    ("SQLite", RGBColor(0x00, 0x3B, 0x57)),
]
start_x = 2.0
for label, color in badges:
    add_pill(slide, label, start_x, 4.55, w=1.65, h=0.38,
             bg=color, fg=WHITE, font_size=12)
    start_x += 1.85

gradient_bar(slide, y=7.42, h=0.08)
add_text(slide, "2025-2026  |  Nagy Kristóf",
         0, 7.1, 13.33, 0.35,
         font_size=11, color=RGBColor(0x9C, 0xA3, 0xAF),
         align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – TARTALOMJEGYZÉK
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, DARK_BG)
gradient_bar(slide, 0)

add_text(slide, "Tartalom", 0.5, 0.18, 12, 0.7,
         font_size=32, bold=True, color=CREAM, align=PP_ALIGN.LEFT)

chapters = [
    ("01", "Projekt bemutatása",         "A NeoCoffee webshop céljai és főbb jellemzői"),
    ("02", "Funkciók áttekintése",        "Vásárlási folyamat, kosár, pénztár, admin"),
    ("03", "Technológiai Stack",          "Frontend, backend, build és tesztelési eszközök"),
    ("04", "Architektúra",               "MVCS minta, rétegek, adatfolyam"),
    ("05", "Adatbázis",                  "SQLite séma, táblák, indexek"),
    ("06", "Backend API",                "Express.js végpontok, JWT autentikáció"),
    ("07", "Frontend komponensek",       "SvelteKit útvonalak, UI komponensek"),
    ("08", "Állapotkezelés",             "Svelte stores, reaktív frissítések"),
    ("09", "Vizuális design",            "Glassmorphism, animációk, témarendszer"),
    ("10", "Tesztelés",                  "Vitest unit és integrációs tesztek"),
    ("11", "Fejlesztési parancsok",      "NPM szkriptek, futtatás, build"),
    ("12", "Összefoglalás",             "Eredmények és lehetséges fejlesztések"),
]

cols = [chapters[:6], chapters[6:]]
col_x = [0.5, 6.9]

for ci, col in enumerate(cols):
    cx = col_x[ci]
    for i, (num, title, desc) in enumerate(col):
        row_y = 1.1 + i * 1.0
        add_pill(slide, num, cx, row_y + 0.08, w=0.52, h=0.4,
                 bg=COFFEE_MID, fg=CREAM, font_size=13)
        add_text(slide, title, cx + 0.65, row_y, 5.2, 0.42,
                 font_size=14, bold=True, color=CREAM)
        add_text(slide, desc, cx + 0.65, row_y + 0.38, 5.2, 0.42,
                 font_size=11, color=LIGHT_GRAY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – PROJEKT BEMUTATÁSA
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
section_header(slide, "01  Projekt bemutatása",
               "NeoCoffee – modern kávé webáruház glassmorphism dizájnnal")

# Bal oldal: leírás
info_lines = [
    "A NeoCoffee egy teljes körű kávé webáruház alkalmazás,",
    "amelyet SvelteKit frontend és Express.js backend épít fel.",
    "",
    "A vásárlók 6 különböző kávét rendelhetnek testreszabható",
    "opciókkal (tej típus, cukor mennyiség), majd online adják",
    "le rendeléseiket – amit az admin felületen kezelnek.",
    "",
    "Különlegesség: a cégünk speciális pohár technológiája",
    "8 órán át melegen tartja a kávét!",
]
add_multiline_text(slide, info_lines, 0.5, 2.35, 6.8, 4.5,
                   font_size=15, color=LIGHT_GRAY)

# Jobb oldal: key facts kártyák
facts = [
    ("6", "Kávétermék"),
    ("3", "Tej opció"),
    ("4", "Főoldal route"),
    ("110+", "Teszt eset"),
]
fx = 8.0
for i, (val, label) in enumerate(facts):
    fy = 2.4 + i * 1.1
    card = slide.shapes.add_shape(9,
        Inches(fx), Inches(fy), Inches(4.7), Inches(0.88))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x22, 0x14, 0x0A)
    card.line.color.rgb = COFFEE_MID
    add_text(slide, val, fx + 0.15, fy + 0.05, 1.1, 0.75,
             font_size=28, bold=True, color=COFFEE_LIGHT, align=PP_ALIGN.RIGHT)
    add_text(slide, label, fx + 1.4, fy + 0.22, 3.0, 0.5,
             font_size=14, bold=False, color=CREAM)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – FUNKCIÓK ÁTTEKINTÉSE
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
section_header(slide, "02  Funkciók áttekintése",
               "A webáruház főbb funkcionális területei")

features = [
    ("🛒", "Termékek & Kosár",
     ["• 6 kávé termék (Cappuccino, Espresso, Ristretto,",
      "  Latte, Americano, Doppio)",
      "• Ár: 650–900 Ft / termék",
      "• Tej típus: Mandula / Tehén / Zab",
      "• Cukor: Nincs / 1 kanál / 2 kanál",
      "• Mennyiség szabályozás (+/–)",
      "• Valós idejű végösszeg"]),
    ("💳", "Pénztár (Checkout)",
     ["• Rendelési űrlap validációval",
      "• Mezők: név, email, telefon,",
      "  irányítószám, település, cím",
      "• Rendelés összesítő oldalsáv",
      "• Automatikus kosár ürítés",
      "• Visszairányítás főoldalra"]),
    ("🔐", "Admin Felület",
     ["• JWT token alapú belépés",
      "  (admin / Minad123!)",
      "• Rendelések táblázatos listája",
      "• Vevő adatok megtekintése",
      "• Postázás státusz kezelés",
      "• Rendelés részletek modal",
      "• Rendelés törlés"]),
    ("🌙", "Téma Rendszer",
     ["• Világos ↔ Sötét mód váltás",
      "• localStorage mentés",
      "• 700ms smooth átmenet",
      "• Dekorációk:",
      "  ☀ Felhők + nap (light)",
      "  ★ Csillagok + hold (dark)"]),
]

cols_x = [0.35, 3.6, 6.85, 10.1]
for i, (icon, title, lines) in enumerate(features):
    cx = cols_x[i]
    # Kártya
    card = slide.shapes.add_shape(9,
        Inches(cx), Inches(2.35), Inches(2.95), Inches(4.6))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x15, 0x0D, 0x06)
    card.line.color.rgb = COFFEE_MID
    # Ikon
    add_text(slide, icon, cx + 0.1, 2.45, 2.75, 0.55,
             font_size=24, align=PP_ALIGN.CENTER, color=CREAM)
    # Cím
    add_text(slide, title, cx + 0.1, 2.95, 2.75, 0.45,
             font_size=13, bold=True, color=COFFEE_LIGHT,
             align=PP_ALIGN.CENTER)
    # Tartalom
    add_multiline_text(slide, lines, cx + 0.15, 3.45, 2.65, 3.3,
                       font_size=11, color=LIGHT_GRAY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – TECHNOLÓGIAI STACK
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
section_header(slide, "03  Technológiai Stack",
               "A projekt által felhasznált főbb technológiák")

categories = [
    ("Frontend", COFFEE_MID, [
        ("SvelteKit 2.49.1",   "Meta-framework (SSR, routing, build)"),
        ("Svelte 5.45.6",      "Reaktív komponens framework (Runes API)"),
        ("Tailwind CSS 4.1.17","Utility-first CSS framework"),
        ("shadcn-svelte 1.1.0","Újrafelhasználható UI komponensek"),
        ("TypeScript 5.9.3",   "Statikus típusellenőrzés"),
        ("@lucide/svelte",     "Ikon könyvtár"),
    ]),
    ("Backend", RGBColor(0x00, 0x3B, 0x57), [
        ("Express.js 4.22.1",  "REST API szerver (Node.js)"),
        ("better-sqlite3",     "Szinkron SQLite driver"),
        ("bcrypt 5.1.1",       "Jelszó hash-elés"),
        ("jsonwebtoken 9.0.3", "JWT token autentikáció"),
        ("cors 2.8.5",         "Cross-Origin Resource Sharing"),
    ]),
    ("Build & Test", RGBColor(0x1A, 0x1A, 0x2E), [
        ("Vite 7.2.6",         "Gyors build tool és dev szerver"),
        ("Vitest 4.0.18",      "Unit és integrációs tesztelés"),
        ("concurrently",       "Frontend + backend párhuzamos futtatás"),
        ("jsdom 27.4.0",       "DOM szimuláció tesztekhez"),
    ]),
]

col_x = [0.35, 4.6, 8.85]
col_w = 4.05

for ci, (cat_title, cat_color, items) in enumerate(categories):
    cx = col_x[ci]
    # Fejléc
    hdr = slide.shapes.add_shape(9,
        Inches(cx), Inches(2.35), Inches(col_w), Inches(0.52))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = cat_color
    hdr.line.fill.background()
    add_text(slide, cat_title, cx + 0.1, 2.38, col_w - 0.2, 0.45,
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Elemek
    for i, (name, desc) in enumerate(items):
        iy = 3.0 + i * 0.68
        add_text(slide, name, cx + 0.2, iy, col_w - 0.3, 0.35,
                 font_size=13, bold=True, color=COFFEE_LIGHT)
        add_text(slide, desc, cx + 0.2, iy + 0.3, col_w - 0.3, 0.32,
                 font_size=10, color=LIGHT_GRAY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – ARCHITEKTÚRA
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
section_header(slide, "04  Architektúra",
               "MVCS (Model-View-Controller-Store) minta – módosított változat")

# Rétegek ábrája
layers = [
    ("VIEW LAYER",   "Svelte komponensek – UI renderelés, user interakciók",  RGBColor(0x4E, 0x27, 0x0D)),
    ("STORE LAYER",  "Svelte stores – reaktív állapot, derived értékek",       RGBColor(0x7B, 0x44, 0x18)),
    ("LOGIC LAYER",  "cart.ts, theme.ts – üzleti logika, side-effects",        RGBColor(0xA0, 0x62, 0x28)),
    ("MODEL LAYER",  "types.ts, products.ts – adatstruktúrák, típusok",        RGBColor(0xC4, 0x8B, 0x4A)),
]

lx = 0.5
for i, (ltitle, ldesc, lcolor) in enumerate(layers):
    ly = 2.35 + i * 0.98
    box = slide.shapes.add_shape(9,
        Inches(lx), Inches(ly), Inches(7.5), Inches(0.82))
    box.fill.solid(); box.fill.fore_color.rgb = lcolor
    box.line.fill.background()
    add_text(slide, ltitle, lx + 0.15, ly + 0.05, 2.5, 0.7,
             font_size=13, bold=True, color=WHITE)
    add_text(slide, ldesc, lx + 2.7, ly + 0.15, 5.0, 0.55,
             font_size=12, color=CREAM)
    if i < 3:
        add_text(slide, "▼", lx + 3.5, ly + 0.78, 0.7, 0.25,
                 font_size=14, bold=True, color=COFFEE_LIGHT,
                 align=PP_ALIGN.CENTER)

# Jobb oldal: adatfolyam szöveg
flow_lines = [
    "Adatfolyam – kosár példán:",
    "",
    "1. Felhasználó beállítja tej/cukor/mennyiség",
    "2. ProductCard meghívja handleAddToCart()",
    "3. addToCart() létrehoz egy CartItem objektumot",
    "4. cartItems.update() → store frissül",
    "5. cartTotal derived store újraszámít",
    "6. Svelte automatikusan frissíti a DOM-ot",
    "7. A végösszeg azonnal megjelenik",
]
flow_colors = [COFFEE_LIGHT] + [LIGHT_GRAY] * (len(flow_lines)-1)
add_multiline_text(slide, flow_lines, 8.3, 2.35, 4.6, 4.5,
                   font_size=13, line_colors=flow_colors)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – ADATBÁZIS
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
section_header(slide, "05  Adatbázis",
               "SQLite 3 – könnyűsúlyú, fájl alapú relációs adatbázis")

tables = [
    ("termekek", COFFEE_BROWN,
     ["id  TEXT  PK", "nev  TEXT  NOT NULL",
      "ar  INTEGER  NOT NULL", "kep_url  TEXT"]),
    ("rendelesek", RGBColor(0x1E, 0x40, 0x3E),
     ["id  TEXT  PK", "vevo_nev / email / telefon",
      "iranyitoszam / telepules / utca_hazszam",
      "megrendelve  TIMESTAMP",
      "postazva  INTEGER  (0/1)",
      "postazva_datum  TEXT"]),
    ("rendeles_tetelek", RGBColor(0x1A, 0x27, 0x47),
     ["id  TEXT  PK",
      "rendeles_id  FK → rendelesek",
      "termek_nev / termek_ar  (snapshot)",
      "mennyiseg  INTEGER",
      "tej / cukor  TEXT"]),
    ("admin", RGBColor(0x3B, 0x1A, 0x47),
     ["id  TEXT  PK",
      "felhasznalonev  TEXT  UNIQUE",
      "jelszo_hash  TEXT  (bcrypt)",
      "letrehozva  TIMESTAMP"]),
]

col_x = [0.35, 3.65, 6.95, 10.2]
col_w = 3.0

for ci, (tname, tcolor, fields) in enumerate(tables):
    cx = col_x[ci]
    # Tábla fejléc
    hdr = slide.shapes.add_shape(9,
        Inches(cx), Inches(2.35), Inches(col_w), Inches(0.5))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = tcolor
    hdr.line.fill.background()
    add_text(slide, tname, cx + 0.1, 2.38, col_w - 0.2, 0.42,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Mezők
    for fi, field in enumerate(fields):
        fy = 2.97 + fi * 0.58
        row_bg = slide.shapes.add_shape(1,
            Inches(cx), Inches(fy), Inches(col_w), Inches(0.52))
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = (
            RGBColor(0x12, 0x0A, 0x04) if fi % 2 == 0
            else RGBColor(0x1C, 0x10, 0x06))
        row_bg.line.fill.background()
        add_text(slide, field, cx + 0.1, fy + 0.05, col_w - 0.2, 0.42,
                 font_size=11, color=LIGHT_GRAY)

# Megjegyzések
add_text(slide, "• ON DELETE CASCADE: rendeles_tetelek törlődnek a rendeléssel együtt",
         0.35, 6.6, 12.6, 0.42, font_size=11, color=LIGHT_GRAY)
add_text(slide, "• Snapshot árazás: az ár a rendeléskor rögzített érték, árváltozástól független",
         0.35, 6.98, 12.6, 0.42, font_size=11, color=LIGHT_GRAY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – BACKEND API
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
section_header(slide, "06  Backend API",
               "Express.js REST API – JWT autentikáció, CORS, SQLite")

# Bal rész: végpontok
endpoints = [
    # (method, path, auth, desc)
    ("GET",    "/api/products",           "–",      "Összes termék listája"),
    ("GET",    "/api/products/:id",       "–",      "Egy termék részletei"),
    ("POST",   "/api/products",           "Admin",  "Új termék létrehozása"),
    ("PUT",    "/api/products/:id",       "Admin",  "Termék módosítása"),
    ("DELETE", "/api/products/:id",       "Admin",  "Termék törlése"),
    ("POST",   "/api/orders",             "–",      "Új rendelés (publikus)"),
    ("GET",    "/api/orders",             "Admin",  "Rendelések listája"),
    ("GET",    "/api/orders/:id",         "Admin",  "Rendelés részletei"),
    ("PATCH",  "/api/orders/:id/ship",    "Admin",  "Postázás státusz"),
    ("DELETE", "/api/orders/:id",         "Admin",  "Rendelés törlése"),
    ("POST",   "/api/admin/login",        "–",      "Bejelentkezés → JWT"),
    ("POST",   "/api/admin/logout",       "–",      "Kijelentkezés"),
    ("GET",    "/api/admin/verify",       "–",      "Token ellenőrzés"),
]

method_colors = {
    "GET":    RGBColor(0x22, 0x8B, 0x22),
    "POST":   RGBColor(0x31, 0x78, 0xC6),
    "PUT":    RGBColor(0xD9, 0x7A, 0x06),
    "PATCH":  RGBColor(0x7B, 0x44, 0x18),
    "DELETE": RGBColor(0xDC, 0x26, 0x26),
}

table_start_y = 2.38
row_h = 0.36

for ri, (method, path, auth, desc) in enumerate(endpoints):
    ry = table_start_y + ri * row_h
    # Sor háttér
    row_bg = slide.shapes.add_shape(1,
        Inches(0.35), Inches(ry), Inches(12.6), Inches(row_h - 0.02))
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = (
        RGBColor(0x14, 0x0C, 0x04) if ri % 2 == 0
        else RGBColor(0x1C, 0x10, 0x06))
    row_bg.line.fill.background()

    # Method badge
    add_pill(slide, method, 0.4, ry + 0.02, w=0.85, h=0.3,
             bg=method_colors.get(method, COFFEE_MID), fg=WHITE, font_size=9)
    # Path
    add_text(slide, path, 1.35, ry + 0.04, 3.6, 0.3,
             font_size=11, bold=True, color=CREAM)
    # Auth
    auth_color = RGBColor(0xFF, 0xC1, 0x07) if auth == "Admin" else LIGHT_GRAY
    add_text(slide, auth, 5.0, ry + 0.04, 1.0, 0.3,
             font_size=10, bold=(auth == "Admin"), color=auth_color,
             align=PP_ALIGN.CENTER)
    # Leírás
    add_text(slide, desc, 6.1, ry + 0.04, 6.7, 0.3,
             font_size=10, color=LIGHT_GRAY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – FRONTEND KOMPONENSEK
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
section_header(slide, "07  Frontend komponensek",
               "SvelteKit útvonalak és komponens hierarchia")

# Hierarchia szöveg (bal oldal)
hierarchy = [
    "📁  src/",
    "  📁  routes/",
    "    📄  +layout.svelte       ← Gyökér layout",
    "    📄  +page.svelte         ← Főoldal (termékek, kosár)",
    "    📁  checkout/",
    "      📄  +page.svelte      ← Rendelési űrlap",
    "    📁  admin/",
    "      📄  +page.svelte      ← Admin panel (JWT védett)",
    "  📁  lib/",
    "    📄  api.ts               ← HTTP kliens (token kezelés)",
    "    📄  cart.ts              ← Kosár store + addToCart()",
    "    📄  orders.ts            ← Rendelések store",
    "    📄  products.ts          ← 6 termék statikus katalógus",
    "    📄  theme.ts             ← isDarkMode store",
    "    📄  types.ts             ← TypeScript interfészek",
    "    📄  utils.ts             ← Validációs függvények",
    "    📁  components/",
    "      📄  ProductCard.svelte      ← Termék kártya",
    "      📄  ThemeDecorations.svelte ← Deko (csillagok/felhők)",
    "      📁  ui/  (shadcn-svelte)    ← 30+ UI komponens",
]
add_multiline_text(slide, hierarchy, 0.35, 2.35, 7.2, 5.0,
                   font_size=11, color=LIGHT_GRAY)

# Jobb oldal: fontos komponensek
comps = [
    ("ProductCard.svelte",
     "Tej/cukor dropdown, mennyiség\n+/– gombok, kosárba gomb"),
    ("ThemeDecorations.svelte",
     "Generált csillagok (dark) és\nfelhők (light) animációval"),
    ("+page.svelte (főoldal)",
     "6x ProductCard grid,\nCheckout bar végösszeggel"),
    ("+page.svelte (admin)",
     "Bejelentkezés, rendelés táblázat,\nDialog modal, Badge státusz"),
]
cx = 7.8
for i, (cname, cdesc) in enumerate(comps):
    cy = 2.4 + i * 1.22
    card = slide.shapes.add_shape(9,
        Inches(cx), Inches(cy), Inches(5.2), Inches(1.05))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x18, 0x0E, 0x05)
    card.line.color.rgb = COFFEE_MID
    add_text(slide, cname, cx + 0.15, cy + 0.05, 4.9, 0.42,
             font_size=13, bold=True, color=COFFEE_LIGHT)
    add_text(slide, cdesc, cx + 0.15, cy + 0.48, 4.9, 0.5,
             font_size=11, color=LIGHT_GRAY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – ÁLLAPOTKEZELÉS
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
section_header(slide, "08  Állapotkezelés",
               "Svelte stores – reaktív, minimális overhead")

stores = [
    ("cartItems", "Writable<CartItem[]>",
     COFFEE_MID,
     ["Kosár tartalmának tárolása",
      "Inicializálás: üres tömb []",
      "Módosítás: update() metódus",
      "Perzisztencia: nincs (session-only)",
      "Export: addToCart() függvény"]),
    ("cartTotal", "Derived<number>",
     RGBColor(0x06, 0x6B, 0x4A),
     ["Kosár végösszeg automatikus számítás",
      "Függ: cartItems store-tól",
      "Csak olvasható (read-only)",
      "Képlet: Σ(ár × mennyiség)",
      "Automatikus újraszámítás"]),
    ("isDarkMode", "Writable<boolean>",
     RGBColor(0x2D, 0x1B, 0x4E),
     ["Aktív téma tárolása",
      "Inicializálás: localStorage / true",
      "Side-effect: DOM classList",
      "Perzisztencia: localStorage",
      "Export: toggleTheme() függvény"]),
]

col_x = [0.35, 4.6, 8.85]
col_w = 3.9

for ci, (sname, stype, scolor, sbullets) in enumerate(stores):
    cx = col_x[ci]
    hdr = slide.shapes.add_shape(9,
        Inches(cx), Inches(2.35), Inches(col_w), Inches(0.72))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = scolor
    hdr.line.fill.background()
    add_text(slide, sname, cx + 0.1, 2.38, col_w - 0.2, 0.38,
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, stype, cx + 0.1, 2.72, col_w - 0.2, 0.3,
             font_size=11, color=CREAM, align=PP_ALIGN.CENTER)

    for bi, bullet in enumerate(sbullets):
        by = 3.22 + bi * 0.65
        add_text(slide, "▸  " + bullet, cx + 0.2, by, col_w - 0.3, 0.55,
                 font_size=12, color=LIGHT_GRAY)

# Reaktivitás megjegyzés
note = ('Svelte $-prefix auto-subscribe: <span>{$cartTotal} Ft</span> – '
        'automatikus DOM frissítés store változáskor, runtime overhead nélkül.')
add_text(slide,
         '💡  Svelte $-prefix auto-subscribe: "{$cartTotal} Ft" – '
         'automatikus DOM frissítés store változáskor, runtime overhead nélkül.',
         0.35, 6.65, 12.6, 0.62,
         font_size=12, color=LIGHT_GRAY, italic=True)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – VIZUÁLIS DESIGN
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
section_header(slide, "09  Vizuális Design",
               "Glassmorphism, animációk, kétféle témamód")

# Bal: Glassmorphism leírás
glass_lines = [
    "Glassmorphism rétegek:",
    "",
    "1. Háttér blur:   backdrop-blur-xl  (24px)",
    "2. Overlay:       bg-white/[0.08]   (8% fehér)",
    "3. Border:        border-white/[0.15]",
    "4. Árnyék:        shadow-[0_8px_32px_rgba(0,0,0,.4)]",
    "5. Hover:         bg-white/[0.15]  (erősebb)",
    "",
    "Animációk:",
    "",
    "• twinkle – csillag villogás (sötét mód)",
    "• float   – felhő lebegés   (világos mód)",
    "• transition-all duration-300 (hover)",
    "• transition-opacity duration-700 (téma)",
]
glass_colors = [COFFEE_LIGHT] + [LIGHT_GRAY] * (len(glass_lines)-1)
add_multiline_text(slide, glass_lines, 0.35, 2.35, 6.0, 4.8,
                   font_size=13, line_colors=glass_colors)

# Jobb: Kétféle téma
theme_x = 6.7
# Sötét kártya
dark_card = slide.shapes.add_shape(9,
    Inches(theme_x), Inches(2.35), Inches(3.0), Inches(2.1))
dark_card.fill.solid(); dark_card.fill.fore_color.rgb = RGBColor(0x0D, 0x0E, 0x1A)
dark_card.line.color.rgb = PURPLE_ACCENT
add_text(slide, "🌙  Sötét mód", theme_x + 0.1, 2.42, 2.8, 0.45,
         font_size=14, bold=True, color=CREAM, align=PP_ALIGN.CENTER)
dark_items = ["Háttér: #0d0e1a", "Akcentus: pink→lila gradiens",
              "Csillagok + hold", "text-white / text-gray-200"]
for i, item in enumerate(dark_items):
    add_text(slide, "★  " + item, theme_x + 0.15, 2.88 + i * 0.35, 2.7, 0.3,
             font_size=11, color=LIGHT_GRAY)

# Világos kártya
light_x = theme_x + 3.3
light_card = slide.shapes.add_shape(9,
    Inches(light_x), Inches(2.35), Inches(3.0), Inches(2.1))
light_card.fill.solid(); light_card.fill.fore_color.rgb = RGBColor(0x47, 0x55, 0x69)
light_card.line.color.rgb = COFFEE_LIGHT
add_text(slide, "☀  Világos mód", light_x + 0.1, 2.42, 2.8, 0.45,
         font_size=14, bold=True, color=CREAM, align=PP_ALIGN.CENTER)
light_items = ["Háttér: slate-700 → slate-500", "Overlay: bg-white/30",
               "Felhők + nap", "text-gray-800"]
for i, item in enumerate(light_items):
    add_text(slide, "☁  " + item, light_x + 0.15, 2.88 + i * 0.35, 2.7, 0.3,
             font_size=11, color=CREAM)

# Responsive rács
add_text(slide, "Responsive rács rendszer:", 6.7, 4.68, 6.4, 0.42,
         font_size=13, bold=True, color=COFFEE_LIGHT)
resp = [
    "• Mobil (< 768px):   1 oszlop   grid-cols-1",
    "• Tablet (768px+):   2 oszlop   md:grid-cols-2",
    "• Desktop (1024px+): 3 oszlop   lg:grid-cols-3",
    "• Max szélesség:  max-w-6xl mx-auto",
]
add_multiline_text(slide, resp, 6.7, 5.1, 6.4, 2.2,
                   font_size=12, color=LIGHT_GRAY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – TESZTELÉS
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
section_header(slide, "10  Tesztelés",
               "Vitest – unit és integrációs tesztek, 110+ teszteset")

test_files = [
    ("tests/utils.test.ts",  "59",  "Validációs függvények",
     COFFEE_MID,
     ["validateEmail – 12 eset",
      "validatePhone – 14 eset",
      "validateZip – 10 eset",
      "validateCity – 12 eset",
      "validateStreet – 11 eset"]),
    ("tests/cart.test.ts",   "15",  "Kosár store",
     RGBColor(0x06, 0x6B, 0x4A),
     ["addToCart – termék hozzáadás",
      "cartTotal – végösszeg számítás",
      "mennyiség kezelés",
      "módosítók (tej, cukor)",
      "több tétel kezelése"]),
    ("tests/api.test.ts",    "16",  "API kliens",
     RGBColor(0x31, 0x78, 0xC6),
     ["Token set/get/clear",
      "authHeaders generálás",
      "login/logout hívás",
      "hiba kezelés",
      "fetch wrapper"]),
    ("tests/server.test.ts", "~20", "Backend integrációs",
     RGBColor(0x2D, 0x1B, 0x4E),
     ["API végpontok tesztelése",
      "JWT autentikáció",
      "Rendelés CRUD",
      "Termék végpontok",
      "Szerver fut közben"]),
]

col_x = [0.35, 3.65, 6.95, 10.2]
col_w = 3.0

for ci, (fname, count, title, fcolor, bullets) in enumerate(test_files):
    cx = col_x[ci]
    # Fejléc
    hdr = slide.shapes.add_shape(9,
        Inches(cx), Inches(2.35), Inches(col_w), Inches(0.52))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = fcolor
    hdr.line.fill.background()
    add_text(slide, title, cx + 0.1, 2.38, col_w - 0.2, 0.42,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Fájlnév
    add_text(slide, fname, cx + 0.1, 3.0, col_w - 0.2, 0.35,
             font_size=10, bold=False, color=LIGHT_GRAY)
    # Szám badge
    add_pill(slide, count + " eset", cx + (col_w - 1.3) / 2, 3.4,
             w=1.3, h=0.4, bg=fcolor, fg=WHITE, font_size=14)
    # Bullet-ek
    for bi, b in enumerate(bullets):
        by = 3.95 + bi * 0.55
        add_text(slide, "• " + b, cx + 0.1, by, col_w - 0.15, 0.48,
                 font_size=11, color=LIGHT_GRAY)

# Futtatás parancsok
add_text(slide, "Futtatás:", 0.35, 6.72, 3.0, 0.35,
         font_size=12, bold=True, color=COFFEE_LIGHT)
cmds = "npm run test   |   npm run test:watch   |   npm run test:server   |   npm run test:all"
add_text(slide, cmds, 0.35, 7.05, 12.6, 0.35,
         font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – FEJLESZTÉSI PARANCSOK & PROJEKT STRUKTÚRA
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
section_header(slide, "11  Fejlesztési parancsok",
               "NPM szkriptek – fejlesztés, adatbázis, tesztelés, build")

cmds_left = [
    ("▶ Fejlesztés", [
        ("npm run dev",        "Frontend dev szerver  (port 5173)"),
        ("npm run server",     "Backend API szerver   (port 3000)"),
        ("npm run dev:all",    "Mindkettő egyszerre   (concurrently)"),
    ]),
    ("🗄 Adatbázis", [
        ("npm run db:init",    "Adatbázis inicializálása"),
        ("npm run db:reset",   "Adatbázis visszaállítása"),
        ("npm run db:check",   "Adatbázis ellenőrzése"),
        ("npm run db:seed",    "Tesztadatok betöltése"),
    ]),
]
cmds_right = [
    ("🧪 Tesztelés", [
        ("npm run test",       "Unit tesztek (vitest run)"),
        ("npm run test:watch", "Tesztek figyelő módban"),
        ("npm run test:server","Backend integrációs tesztek"),
        ("npm run test:all",   "Összes teszt egyszerre"),
    ]),
    ("📦 Build", [
        ("npm run build",      "Produkciós build (Vite)"),
        ("npm run preview",    "Build előnézet"),
        ("npm run check",      "TypeScript + svelte-check"),
    ]),
]

for gi, group_list in enumerate([cmds_left, cmds_right]):
    base_x = 0.35 + gi * 6.6
    current_y = 2.38
    for gname, gcmds in group_list:
        add_text(slide, gname, base_x, current_y, 6.2, 0.42,
                 font_size=14, bold=True, color=COFFEE_LIGHT)
        current_y += 0.42
        for cmd, desc in gcmds:
            row = slide.shapes.add_shape(1,
                Inches(base_x), Inches(current_y),
                Inches(6.2), Inches(0.48))
            row.fill.solid()
            row.fill.fore_color.rgb = RGBColor(0x14, 0x0C, 0x04)
            row.line.fill.background()
            add_text(slide, cmd, base_x + 0.15, current_y + 0.05,
                     2.4, 0.38, font_size=12, bold=True, color=CREAM)
            add_text(slide, desc, base_x + 2.65, current_y + 0.05,
                     3.4, 0.38, font_size=11, color=LIGHT_GRAY)
            current_y += 0.5
        current_y += 0.2

# URL-ek
add_text(slide, "🌐  Frontend: http://localhost:5173   |   🔌  Backend API: http://localhost:3000",
         0.35, 7.05, 12.6, 0.38, font_size=12, color=CREAM, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 14 – ÖSSZEFOGLALÁS
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, DARK_BG)

# Dekor körök
c1 = slide.shapes.add_shape(9, Inches(-1), Inches(4), Inches(3.5), Inches(3.5))
c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x2D, 0x1A, 0x08)
c1.line.fill.background()
c2 = slide.shapes.add_shape(9, Inches(11), Inches(-1), Inches(3), Inches(3))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x1E, 0x12, 0x30)
c2.line.fill.background()

gradient_bar(slide, y=0.0, h=0.06)

add_text(slide, "12  Összefoglalás",
         0.5, 0.18, 12, 0.62,
         font_size=28, bold=True, color=CREAM, align=PP_ALIGN.LEFT)

achievements = [
    ("✅", "Full-stack alkalmazás",
     "SvelteKit frontend + Express.js backend + SQLite adatbázis"),
    ("✅", "Teljes vásárlási folyamat",
     "Termék kiválasztás → kosár → checkout → rendelés mentés"),
    ("✅", "Admin panel",
     "JWT autentikáció, rendelés kezelés, státusz módosítás"),
    ("✅", "110+ teszt eset",
     "Vitest unit tesztek (utils, cart, api) + backend integrációs tesztek"),
    ("✅", "Glassmorphism design",
     "Animált dekorációk, kétféle téma, reszponzív rács"),
    ("✅", "Reaktív állapotkezelés",
     "Svelte stores – automatikus DOM frissítés, localStorage perzisztencia"),
]

for i, (icon, title, desc) in enumerate(achievements):
    iy = 1.05 + i * 0.85
    add_text(slide, icon, 0.5, iy + 0.18, 0.55, 0.5,
             font_size=16, align=PP_ALIGN.CENTER, color=GREEN_OK)
    add_text(slide, title, 1.1, iy + 0.1, 3.8, 0.42,
             font_size=14, bold=True, color=CREAM)
    add_text(slide, desc, 5.0, iy + 0.1, 8.0, 0.42,
             font_size=13, color=LIGHT_GRAY)

# Jövőbeli fejlesztések
future_lines = [
    "Lehetséges fejlesztések:",
    "• Kosár localStorage perzisztencia",
    "• Felhasználói fiók regisztráció/bejelentkezés",
    "• Email értesítések rendelés után",
    "• Termékkeresés és szűrés",
    "• Fizetési gateway integráció (Stripe)",
]
future_colors = [COFFEE_LIGHT] + [LIGHT_GRAY] * 5

add_multiline_text(slide, future_lines, 0.5, 6.15, 12.5, 1.15,
                   font_size=12, line_colors=future_colors)

gradient_bar(slide, y=7.42, h=0.08)
add_text(slide, "NeoCoffee Webáruház  |  2025-2026  |  Nagy Kristóf",
         0, 7.1, 13.33, 0.35,
         font_size=11, color=RGBColor(0x9C, 0xA3, 0xAF),
         align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# KÓD DIÁK – SEGÉDFÜGGVÉNY
# ═════════════════════════════════════════════════════════════════════════════

# Syntax-highlight színek
SX_KEYWORD  = RGBColor(0xC7, 0x8A, 0xFF)   # lila  – kulcsszavak (import, export, const…)
SX_STRING   = RGBColor(0x9E, 0xCE, 0x6A)   # zöld  – string literálok
SX_COMMENT  = RGBColor(0x6A, 0x73, 0x7D)   # szürke – megjegyzések
SX_TYPE     = RGBColor(0x56, 0xB6, 0xC2)   # kék   – TypeScript típusok
SX_NUMBER   = RGBColor(0xD1, 0x9A, 0x66)   # narancs – számok
SX_FUNC     = RGBColor(0x61, 0xAF, 0xEF)   # világos-kék – függvénynevek
SX_PLAIN    = RGBColor(0xAB, 0xB2, 0xBF)   # alap szöveg
SX_HTML_TAG = RGBColor(0xE0, 0x6C, 0x75)   # piros – HTML/Svelte tagek
SX_ATTR     = RGBColor(0xD1, 0x9A, 0x66)   # narancs – attribútumok
SX_SQL_KW   = RGBColor(0x56, 0xB6, 0xC2)   # kék  – SQL kulcsszavak

CODE_BG     = RGBColor(0x0A, 0x0D, 0x16)   # kód háttér
CODE_BORDER = RGBColor(0x2A, 0x2D, 0x3A)


def code_slide(title, subtitle, filename, code_blocks):
    """
    Kód diavetítő.
    title      – dia fejléc
    subtitle   – kis alcím
    filename   – fájlnév badge (pl. 'src/lib/cart.ts')
    code_blocks – lista: (x, y, w, h, lines, font_size)
                  lines = [(text, color), ...]
    """
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, DARK_BG)
    gradient_bar(slide, 0)

    add_text(slide, title, 0.35, 0.1, 9.0, 0.55,
             font_size=22, bold=True, color=CREAM)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.6, 9.5, 0.4,
                 font_size=12, color=LIGHT_GRAY)
    # Fájlnév badge
    add_pill(slide, filename, 9.8, 0.15, w=3.25, h=0.38,
             bg=CODE_BORDER, fg=RGBColor(0x79, 0xC0, 0xFF), font_size=11)

    for (bx, by, bw, bh, lines, fs) in code_blocks:
        # Kód blokk háttere
        bg = slide.shapes.add_shape(9,
            Inches(bx), Inches(by), Inches(bw), Inches(bh))
        bg.fill.solid(); bg.fill.fore_color.rgb = CODE_BG
        bg.line.color.rgb = CODE_BORDER

        # Kód szöveg
        txb = slide.shapes.add_textbox(
            Inches(bx + 0.18), Inches(by + 0.15),
            Inches(bw - 0.3), Inches(bh - 0.25))
        txb.word_wrap = False
        tf = txb.text_frame
        tf.word_wrap = False
        for li, (text, color) in enumerate(lines):
            para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            para.alignment = PP_ALIGN.LEFT
            run = para.add_run()
            run.text = text
            run.font.size = Pt(fs)
            run.font.color.rgb = color
            run.font.name = "Consolas"
    return slide


def code_line(text, color=SX_PLAIN):
    return (text, color)


def CL(text, color=SX_PLAIN):
    """Shorthand for code_line."""
    return (text, color)


# ─────────────────────────────────────────────────────────────────────────────
# KÓD DIA 1 – TypeScript Típusok (types.ts)
# ─────────────────────────────────────────────────────────────────────────────
code_slide(
    "Kód – TypeScript Típusok",
    "Az adatstruktúrák definiálása – minden réteg alapja  (types.ts)",
    "src/lib/types.ts",
    [
        (0.3, 1.05, 6.2, 6.2,
         [
             CL("// Termék interfész", SX_COMMENT),
             CL("export interface Product {", SX_KEYWORD),
             CL("    id:     string;",        SX_PLAIN),
             CL("    name:   string;",        SX_PLAIN),
             CL("    price:  number;",        SX_PLAIN),
             CL("    image?: string | null;", SX_PLAIN),
             CL("}",                          SX_KEYWORD),
             CL("",                           SX_PLAIN),
             CL("// Módosítók (tej, cukor)", SX_COMMENT),
             CL("export interface Modifiers {", SX_KEYWORD),
             CL('    milk:  string;   // "Mandula tej" | "Tehéntej" | …', SX_PLAIN),
             CL('    sugar: string;   // "Nincs" | "1 kanál" | "2 kanál"', SX_PLAIN),
             CL("}",                          SX_KEYWORD),
             CL("",                           SX_PLAIN),
             CL("// Kosár elem – kibővíti a Product-ot", SX_COMMENT),
             CL("export interface CartItem extends Product {", SX_KEYWORD),
             CL("    quantity:  number;",     SX_PLAIN),
             CL("    modifiers: Modifiers;",  SX_TYPE),
             CL("}",                          SX_KEYWORD),
         ], 11),
        (6.7, 1.05, 6.3, 6.2,
         [
             CL("// Rendelési tétel", SX_COMMENT),
             CL("export interface OrderItem {", SX_KEYWORD),
             CL("    name:      string;",    SX_PLAIN),
             CL("    quantity:  number;",    SX_PLAIN),
             CL("    price:     number;",    SX_PLAIN),
             CL("    modifiers: Modifiers;", SX_TYPE),
             CL("}",                         SX_KEYWORD),
             CL("",                          SX_PLAIN),
             CL("// Teljes rendelés (admin felület)", SX_COMMENT),
             CL("export interface Order {",  SX_KEYWORD),
             CL("    id:          string;",  SX_PLAIN),
             CL("    customerName: string;", SX_PLAIN),
             CL("    email:       string;",  SX_PLAIN),
             CL("    phone:       string;",  SX_PLAIN),
             CL("    address:     Address;", SX_TYPE),
             CL("    orderDate:   string;",  SX_PLAIN),
             CL("    shipped:     boolean;", SX_PLAIN),
             CL("    shippedDate: string | null;", SX_PLAIN),
             CL("    items:       OrderItem[];", SX_TYPE),
             CL("}",                          SX_KEYWORD),
         ], 11),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# KÓD DIA 2 – Kosár Store (cart.ts)
# ─────────────────────────────────────────────────────────────────────────────
code_slide(
    "Kód – Kosár Store",
    "Svelte writable + derived stores, addToCart függvény  (cart.ts)",
    "src/lib/cart.ts",
    [
        (0.3, 1.05, 12.7, 6.2,
         [
             CL("import { writable, derived } from 'svelte/store';", SX_KEYWORD),
             CL("import type { CartItem, Product, Modifiers } from './types';", SX_KEYWORD),
             CL("",                       SX_PLAIN),
             CL("// Írható store – a kosár tartalma", SX_COMMENT),
             CL("export const cartItems = writable<CartItem[]>([]);", SX_FUNC),
             CL("",                       SX_PLAIN),
             CL("// Derived store – automatikusan újraszámítja a végösszeget", SX_COMMENT),
             CL("export const cartTotal = derived(cartItems, ($items) => {", SX_FUNC),
             CL("    return $items.reduce(", SX_PLAIN),
             CL("        (total, item) => total + (item.price * item.quantity),", SX_PLAIN),
             CL("        0",               SX_NUMBER),
             CL("    );",                  SX_PLAIN),
             CL("});",                     SX_KEYWORD),
             CL("",                        SX_PLAIN),
             CL("// Termék hozzáadása a kosárhoz", SX_COMMENT),
             CL("export const addToCart = (", SX_FUNC),
             CL("    product:   Product,",   SX_PLAIN),
             CL("    quantity:  number,",    SX_PLAIN),
             CL("    modifiers: Modifiers",  SX_TYPE),
             CL(") => {",                    SX_KEYWORD),
             CL("    cartItems.update(items => {", SX_FUNC),
             CL("        const newItem: CartItem = { ...product, quantity, modifiers };", SX_PLAIN),
             CL("        return [...items, newItem];  // új elem hozzáfűzése", SX_PLAIN),
             CL("    });",                   SX_PLAIN),
             CL("};",                        SX_KEYWORD),
         ], 11),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# KÓD DIA 3 – Téma Store (theme.ts)
# ─────────────────────────────────────────────────────────────────────────────
code_slide(
    "Kód – Téma Store",
    "isDarkMode store – localStorage szinkron, DOM osztály kezelés  (theme.ts)",
    "src/lib/theme.ts",
    [
        (0.3, 1.05, 12.7, 6.2,
         [
             CL("import { writable } from 'svelte/store';", SX_KEYWORD),
             CL("import { browser } from '$app/environment';", SX_KEYWORD),
             CL("",                          SX_PLAIN),
             CL("// Alapérték: localStorage-ból, szerveren mindig dark", SX_COMMENT),
             CL("const defaultValue = browser", SX_PLAIN),
             CL("    ? (localStorage.getItem('theme') !== 'light')", SX_STRING),
             CL("    : true;",                SX_KEYWORD),
             CL("",                          SX_PLAIN),
             CL("export const isDarkMode = writable<boolean>(defaultValue);", SX_FUNC),
             CL("",                          SX_PLAIN),
             CL("if (browser) {",            SX_KEYWORD),
             CL("    // Azonnali DOM init – megelőzi a téma-villanást", SX_COMMENT),
             CL("    if (defaultValue) {",    SX_KEYWORD),
             CL("        document.documentElement.classList.add('dark');", SX_PLAIN),
             CL("    } else {",              SX_KEYWORD),
             CL("        document.documentElement.classList.remove('dark');", SX_PLAIN),
             CL("    }",                     SX_KEYWORD),
             CL("",                          SX_PLAIN),
             CL("    // Store változás → DOM + localStorage szinkronizálás", SX_COMMENT),
             CL("    isDarkMode.subscribe((value) => {", SX_FUNC),
             CL("        if (value) {",       SX_KEYWORD),
             CL("            document.documentElement.classList.add('dark');", SX_PLAIN),
             CL("            localStorage.setItem('theme', 'dark');", SX_STRING),
             CL("        } else {",           SX_KEYWORD),
             CL("            document.documentElement.classList.remove('dark');", SX_PLAIN),
             CL("            localStorage.setItem('theme', 'light');", SX_STRING),
             CL("        }",                  SX_KEYWORD),
             CL("    });",                    SX_PLAIN),
             CL("}",                          SX_KEYWORD),
         ], 11),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# KÓD DIA 4 – Validációs Függvények (utils.ts)
# ─────────────────────────────────────────────────────────────────────────────
code_slide(
    "Kód – Validációs Függvények",
    "Checkout form mezők validálása regex-szel  (utils.ts)",
    "src/lib/utils.ts",
    [
        (0.3, 1.05, 6.2, 6.2,
         [
             CL("// Email validáció", SX_COMMENT),
             CL("export function validateEmail(email: string): boolean {", SX_FUNC),
             CL("    const re = /^[^\\s@]+@[^\\s@]+\\.[^\\s@]+$/;", SX_STRING),
             CL("    return re.test(email.trim());", SX_PLAIN),
             CL("}",                   SX_KEYWORD),
             CL("",                    SX_PLAIN),
             CL("// Magyar telefonszám (+36 / 06 + 20|30|70 + 7 szám)", SX_COMMENT),
             CL("export function validatePhone(phone: string): boolean {", SX_FUNC),
             CL("    const clean = phone.replace(/[\\s\\-]/g, '');", SX_PLAIN),
             CL("    const re = /^(\\+36|06)(20|30|70)\\d{7}$/;", SX_STRING),
             CL("    return re.test(clean);", SX_PLAIN),
             CL("}",                   SX_KEYWORD),
             CL("",                    SX_PLAIN),
             CL("// Irányítószám – pontosan 4 számjegy", SX_COMMENT),
             CL("export function validateZip(zip: string): boolean {", SX_FUNC),
             CL("    return /^\\d{4}$/.test(zip.trim());", SX_STRING),
             CL("}",                   SX_KEYWORD),
             CL("",                    SX_PLAIN),
             CL("// Település – magyar betűk, min. 2 karakter", SX_COMMENT),
             CL("export function validateCity(city: string): boolean {", SX_FUNC),
             CL("    const re = /^[a-záéíóöőúüűA-ZÁÉÍÓÖŐÚÜŰ\\s\\-]{2,}$/;", SX_STRING),
             CL("    return re.test(city.trim());", SX_PLAIN),
             CL("}",                   SX_KEYWORD),
         ], 10.5),
        (6.7, 1.05, 6.3, 6.2,
         [
             CL("// Utca, házszám – betűk, számok, pont, kötőjel", SX_COMMENT),
             CL("export function validateStreet(street: string): boolean {", SX_FUNC),
             CL("    const re = /^[a-záéíóöőúüűA-Z0-9\\s\\.\\-,]{3,}$/;", SX_STRING),
             CL("    return re.test(street.trim());", SX_PLAIN),
             CL("}",                    SX_KEYWORD),
             CL("",                     SX_PLAIN),
             CL("// Hibaüzenet konstansok", SX_COMMENT),
             CL("export const validationErrors = {", SX_KEYWORD),
             CL("  email:",             SX_PLAIN),
             CL("    'Helyes formátum: pelda@email.hu',", SX_STRING),
             CL("  phone:",             SX_PLAIN),
             CL("    'Magyar szám: +36 30 123 4567',", SX_STRING),
             CL("  zip:",               SX_PLAIN),
             CL("    '4 számjegyű irányítószám',",    SX_STRING),
             CL("  city:",              SX_PLAIN),
             CL("    'Min. 2 karakter, csak betűk',", SX_STRING),
             CL("  street:",            SX_PLAIN),
             CL("    'Min. 3 karakter',",             SX_STRING),
             CL("  required: 'Kötelező mező'",        SX_STRING),
             CL("};",                   SX_KEYWORD),
         ], 10.5),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# KÓD DIA 5 – ProductCard komponens (Svelte 5 Runes)
# ─────────────────────────────────────────────────────────────────────────────
code_slide(
    "Kód – ProductCard Komponens",
    "Svelte 5 Runes API: $state, $props – termék kártya  (ProductCard.svelte)",
    "src/lib/components/ProductCard.svelte",
    [
        (0.3, 1.05, 6.2, 6.2,
         [
             CL("<script lang=\"ts\">",       SX_HTML_TAG),
             CL("  import { addToCart }",     SX_KEYWORD),
             CL("    from '$lib/cart';",       SX_STRING),
             CL("  import type { Product }",  SX_KEYWORD),
             CL("    from '$lib/types';",      SX_STRING),
             CL("  import * as Select",        SX_KEYWORD),
             CL("    from '$lib/components/ui/select';", SX_STRING),
             CL("",                            SX_PLAIN),
             CL("  // Svelte 5 $props rune",  SX_COMMENT),
             CL("  let { product }: { product: Product }", SX_PLAIN),
             CL("      = $props();",           SX_FUNC),
             CL("",                            SX_PLAIN),
             CL("  // Svelte 5 $state rune",  SX_COMMENT),
             CL("  let quantity = $state(1);", SX_FUNC),
             CL("  let milk     = $state('Nincs');", SX_FUNC),
             CL("  let sugar    = $state('Nincs');", SX_FUNC),
             CL("",                            SX_PLAIN),
             CL("  function decrease() {",     SX_FUNC),
             CL("    if (quantity > 1) quantity--;", SX_PLAIN),
             CL("  }",                          SX_KEYWORD),
             CL("  function increase() { quantity++; }", SX_FUNC),
             CL("",                             SX_PLAIN),
             CL("  function handleAddToCart() {", SX_FUNC),
             CL("    addToCart(product, quantity,", SX_PLAIN),
             CL("      { milk, sugar });",       SX_PLAIN),
             CL("    quantity = 1; // reset",   SX_COMMENT),
             CL("  }",                           SX_KEYWORD),
             CL("</script>",                     SX_HTML_TAG),
         ], 10.5),
        (6.7, 1.05, 6.3, 6.2,
         [
             CL("<!-- Termék kép vagy ikon -->", SX_COMMENT),
             CL("{#if product.image}",           SX_HTML_TAG),
             CL("  <img src={product.image}",    SX_ATTR),
             CL("    alt={product.name} />",      SX_ATTR),
             CL("{:else}",                        SX_HTML_TAG),
             CL("  <Coffee />",                   SX_HTML_TAG),
             CL("{/if}",                          SX_HTML_TAG),
             CL("",                               SX_PLAIN),
             CL("<!-- Tej típus választó -->",    SX_COMMENT),
             CL("<Select.Root type=\"single\"",   SX_HTML_TAG),
             CL("  bind:value={milk as any}>",    SX_ATTR),
             CL("  <Select.Item value=\"Mandula tej\" />", SX_PLAIN),
             CL("  <Select.Item value=\"Tehéntej\" />",    SX_PLAIN),
             CL("  <Select.Item value=\"Zabtej\" />",      SX_PLAIN),
             CL("</Select.Root>",                 SX_HTML_TAG),
             CL("",                               SX_PLAIN),
             CL("<!-- Mennyiség szabályozás -->", SX_COMMENT),
             CL("<Button onclick={decrease}>",    SX_HTML_TAG),
             CL("  <Minus size={14} />",          SX_PLAIN),
             CL("</Button>",                      SX_HTML_TAG),
             CL("<span>{quantity} db</span>",     SX_PLAIN),
             CL("<Button onclick={increase}>",    SX_HTML_TAG),
             CL("  <Plus size={14} />",           SX_PLAIN),
             CL("</Button>",                      SX_HTML_TAG),
             CL("",                               SX_PLAIN),
             CL("<Button onclick={handleAddToCart}>", SX_HTML_TAG),
             CL("  Kosárba",                      SX_STRING),
             CL("</Button>",                      SX_HTML_TAG),
         ], 10.5),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# KÓD DIA 6 – Főoldal (+page.svelte)
# ─────────────────────────────────────────────────────────────────────────────
code_slide(
    "Kód – Főoldal",
    "Termék grid, checkout bar, témaváltó gomb  (+page.svelte)",
    "src/routes/+page.svelte",
    [
        (0.3, 1.05, 6.2, 6.2,
         [
             CL("<script lang=\"ts\">",        SX_HTML_TAG),
             CL("  import { cartTotal }",      SX_KEYWORD),
             CL("    from '$lib/cart';",        SX_STRING),
             CL("  import { isDarkMode }",      SX_KEYWORD),
             CL("    from '$lib/theme';",       SX_STRING),
             CL("  import ProductCard",         SX_KEYWORD),
             CL("    from '$lib/components/ProductCard.svelte';", SX_STRING),
             CL("",                             SX_PLAIN),
             CL("  let { data } = $props();",   SX_FUNC),
             CL("",                             SX_PLAIN),
             CL("  function toggleTheme() {",   SX_FUNC),
             CL("    isDarkMode.update(v => !v);", SX_PLAIN),
             CL("  }",                           SX_KEYWORD),
             CL("</script>",                    SX_HTML_TAG),
             CL("",                             SX_PLAIN),
             CL("<!-- Téma váltó + Admin gomb -->", SX_COMMENT),
             CL("<Button onclick={toggleTheme}>", SX_HTML_TAG),
             CL("  {#if $isDarkMode}",           SX_HTML_TAG),
             CL("    <Sun />",                   SX_PLAIN),
             CL("  {:else}",                     SX_HTML_TAG),
             CL("    <Moon />",                  SX_PLAIN),
             CL("  {/if}",                       SX_HTML_TAG),
             CL("</Button>",                     SX_HTML_TAG),
         ], 10.5),
        (6.7, 1.05, 6.3, 6.2,
         [
             CL("<!-- Termék rács (3 oszlop desktop) -->", SX_COMMENT),
             CL("<div class=\"grid grid-cols-1", SX_ATTR),
             CL("  md:grid-cols-2",              SX_ATTR),
             CL("  lg:grid-cols-3 gap-6\">",     SX_ATTR),
             CL("",                              SX_PLAIN),
             CL("  {#each data.products as product (product.id)}", SX_HTML_TAG),
             CL("    <ProductCard {product} />", SX_HTML_TAG),
             CL("  {/each}",                     SX_HTML_TAG),
             CL("</div>",                        SX_HTML_TAG),
             CL("",                              SX_PLAIN),
             CL("<!-- Checkout sáv – végösszeg -->", SX_COMMENT),
             CL("<Card.Root class=\"w-full max-w-lg\">", SX_HTML_TAG),
             CL("  <div class=\"flex justify-between\">", SX_HTML_TAG),
             CL("    <span>Összesen:</span>",    SX_PLAIN),
             CL("    <span>{$cartTotal} Ft</span>", SX_FUNC),
             CL("  </div>",                      SX_HTML_TAG),
             CL("  <Button href=\"/checkout\">", SX_HTML_TAG),
             CL("    <ShoppingCart />",           SX_PLAIN),
             CL("    Megrendelés",               SX_STRING),
             CL("  </Button>",                   SX_HTML_TAG),
             CL("</Card.Root>",                  SX_HTML_TAG),
         ], 10.5),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# KÓD DIA 7 – Checkout oldal ($effect, validáció)
# ─────────────────────────────────────────────────────────────────────────────
code_slide(
    "Kód – Checkout Oldal",
    "Svelte 5 $state, $effect, form validáció, rendelés küldés  (checkout/+page.svelte)",
    "src/routes/checkout/+page.svelte",
    [
        (0.3, 1.05, 6.2, 6.2,
         [
             CL("// Svelte 5 $state – form mezők", SX_COMMENT),
             CL("let fullName = $state('');", SX_FUNC),
             CL("let email    = $state('');", SX_FUNC),
             CL("let phone    = $state('');", SX_FUNC),
             CL("let zip      = $state('');", SX_FUNC),
             CL("let city     = $state('');", SX_FUNC),
             CL("let street   = $state('');", SX_FUNC),
             CL("",                           SX_PLAIN),
             CL("// $effect – form adatok auto-mentése", SX_COMMENT),
             CL("$effect(() => {",            SX_FUNC),
             CL("  if (!saveFormData) return;", SX_PLAIN),
             CL("  localStorage.setItem(",   SX_PLAIN),
             CL("    CHECKOUT_STORAGE_KEY,",  SX_STRING),
             CL("    JSON.stringify({",       SX_PLAIN),
             CL("      fullName, email,",     SX_PLAIN),
             CL("      phone, zip,",          SX_PLAIN),
             CL("      city, street",         SX_PLAIN),
             CL("    })",                     SX_PLAIN),
             CL("  );",                       SX_PLAIN),
             CL("});",                        SX_KEYWORD),
             CL("",                           SX_PLAIN),
             CL("// onMount – adatok visszatöltése", SX_COMMENT),
             CL("onMount(() => {",            SX_FUNC),
             CL("  const saved = localStorage",    SX_PLAIN),
             CL("    .getItem(CHECKOUT_STORAGE_KEY);", SX_STRING),
             CL("  if (saved) {",             SX_KEYWORD),
             CL("    const d = JSON.parse(saved);", SX_PLAIN),
             CL("    fullName = d.fullName || '';", SX_PLAIN),
             CL("    email    = d.email    || '';", SX_PLAIN),
             CL("  }",                        SX_KEYWORD),
             CL("});",                        SX_KEYWORD),
         ], 10.5),
        (6.7, 1.05, 6.3, 6.2,
         [
             CL("// Rendelés leadása", SX_COMMENT),
             CL("async function placeOrder() {", SX_FUNC),
             CL("  // 1. Form validáció",      SX_COMMENT),
             CL("  if (!validateEmail(email)) {", SX_PLAIN),
             CL("    emailError = validationErrors.email;", SX_PLAIN),
             CL("    return;",                  SX_KEYWORD),
             CL("  }",                          SX_KEYWORD),
             CL("  if (!validatePhone(phone)) {", SX_PLAIN),
             CL("    phoneError = validationErrors.phone;", SX_PLAIN),
             CL("    return;",                  SX_KEYWORD),
             CL("  }",                          SX_KEYWORD),
             CL("  // … zip, city, street …",  SX_COMMENT),
             CL("",                             SX_PLAIN),
             CL("  // 2. API hívás",            SX_COMMENT),
             CL("  const order = await createOrder({", SX_FUNC),
             CL("    vevo_nev:     fullName,",  SX_PLAIN),
             CL("    email,  telefon: phone,",  SX_PLAIN),
             CL("    iranyitoszam: zip,",        SX_PLAIN),
             CL("    telepules:    city,",       SX_PLAIN),
             CL("    utca_hazszam: street,",     SX_PLAIN),
             CL("    items: $cartItems.map(i =>({ …i }))", SX_PLAIN),
             CL("  });",                         SX_KEYWORD),
             CL("",                              SX_PLAIN),
             CL("  // 3. Kosár ürítése + redirect", SX_COMMENT),
             CL("  cartItems.set([]);",          SX_FUNC),
             CL("  goto('/');",                  SX_PLAIN),
             CL("}",                             SX_KEYWORD),
         ], 10.5),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# KÓD DIA 8 – Admin oldal ($derived, async API)
# ─────────────────────────────────────────────────────────────────────────────
code_slide(
    "Kód – Admin Felület",
    "JWT login, rendelés táblázat, $derived szűrő  (admin/+page.svelte)",
    "src/routes/admin/+page.svelte",
    [
        (0.3, 1.05, 6.2, 6.2,
         [
             CL("// Svelte 5 $state – állapotok", SX_COMMENT),
             CL("let isLoggedIn  = $state(false);", SX_FUNC),
             CL("let orders      = $state<ApiOrder[]>([]);", SX_FUNC),
             CL("let statusFilter= $state<string>('all');", SX_FUNC),
             CL("",                               SX_PLAIN),
             CL("// $derived – szűrt rendelések", SX_COMMENT),
             CL("let filteredOrders = $derived(",  SX_FUNC),
             CL("  statusFilter === 'all'",         SX_PLAIN),
             CL("    ? orders",                     SX_PLAIN),
             CL("    : statusFilter === 'pending'", SX_PLAIN),
             CL("      ? orders.filter(o =>",       SX_PLAIN),
             CL("          o.postazva === 0)",       SX_NUMBER),
             CL("      : orders.filter(o =>",        SX_PLAIN),
             CL("          o.postazva === 1)",        SX_NUMBER),
             CL(");",                               SX_KEYWORD),
             CL("",                                SX_PLAIN),
             CL("// Bejelentkezés", SX_COMMENT),
             CL("async function handleLogin() {",  SX_FUNC),
             CL("  await api.login(username, password);", SX_PLAIN),
             CL("  isLoggedIn = true;",             SX_PLAIN),
             CL("  await loadOrders();",            SX_PLAIN),
             CL("}",                                SX_KEYWORD),
         ], 10.5),
        (6.7, 1.05, 6.3, 6.2,
         [
             CL("// Postázás státusz váltás", SX_COMMENT),
             CL("async function toggleShipped(", SX_FUNC),
             CL("    orderId: string,",           SX_PLAIN),
             CL("    currentStatus: number",       SX_PLAIN),
             CL(") {",                             SX_KEYWORD),
             CL("  const newStatus = currentStatus === 1 ? 0 : 1;", SX_PLAIN),
             CL("  const result = await api.shipOrder(", SX_FUNC),
             CL("    orderId, newStatus as 0 | 1", SX_TYPE),
             CL("  );",                            SX_PLAIN),
             CL("  // Csak az érintett sort frissítjük", SX_COMMENT),
             CL("  orders = orders.map(o => {",    SX_FUNC),
             CL("    if (o.id !== orderId) return o;", SX_PLAIN),
             CL("    return { ...o,",              SX_PLAIN),
             CL("      postazva: result.postazva,", SX_PLAIN),
             CL("      postazva_datum: result.postazva_datum", SX_PLAIN),
             CL("    };",                          SX_PLAIN),
             CL("  });",                           SX_PLAIN),
             CL("}",                               SX_KEYWORD),
             CL("",                               SX_PLAIN),
             CL("<!-- Táblázat sor – Badge státusz -->", SX_COMMENT),
             CL("{#if order.postazva === 1}",       SX_HTML_TAG),
             CL("  <Badge variant=\"default\">",    SX_HTML_TAG),
             CL("    Postázva",                     SX_STRING),
             CL("  </Badge>",                       SX_HTML_TAG),
             CL("{:else}",                          SX_HTML_TAG),
             CL("  <Badge variant=\"secondary\">",  SX_HTML_TAG),
             CL("    Feldolgozás alatt",             SX_STRING),
             CL("  </Badge>",                        SX_HTML_TAG),
             CL("{/if}",                             SX_HTML_TAG),
         ], 10.5),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# KÓD DIA 9 – Express Server (autentikáció, middleware)
# ─────────────────────────────────────────────────────────────────────────────
code_slide(
    "Kód – Express Szerver: Autentikáció",
    "JWT generálás, authMiddleware, bejelentkezés végpont  (server.cjs)",
    "server.cjs",
    [
        (0.3, 1.05, 6.2, 6.2,
         [
             CL("// JWT konfiguráció", SX_COMMENT),
             CL("const JWT_SECRET =",            SX_PLAIN),
             CL("  process.env.JWT_SECRET ||",   SX_KEYWORD),
             CL("  'neocoffee-secret-key';",     SX_STRING),
             CL("const JWT_EXPIRES_IN = '24h';", SX_STRING),
             CL("",                              SX_PLAIN),
             CL("// Token generálás", SX_COMMENT),
             CL("function generateToken(admin) {", SX_FUNC),
             CL("  return jwt.sign(",             SX_PLAIN),
             CL("    { id: admin.id,",            SX_PLAIN),
             CL("      felhasznalonev: admin.felhasznalonev },", SX_PLAIN),
             CL("    JWT_SECRET,",                SX_STRING),
             CL("    { expiresIn: JWT_EXPIRES_IN }", SX_PLAIN),
             CL("  );",                           SX_PLAIN),
             CL("}",                              SX_KEYWORD),
             CL("",                              SX_PLAIN),
             CL("// Auth middleware", SX_COMMENT),
             CL("function authMiddleware(req, res, next) {", SX_FUNC),
             CL("  const auth = req.headers.authorization;", SX_PLAIN),
             CL("  if (!auth || !auth.startsWith('Bearer '))",SX_PLAIN),
             CL("    return res.status(401)",     SX_PLAIN),
             CL("      .json({ error: 'Nincs jogosultság' });", SX_STRING),
             CL("  try {",                        SX_KEYWORD),
             CL("    req.admin = jwt.verify(",    SX_FUNC),
             CL("      auth.split(' ')[1], JWT_SECRET);", SX_STRING),
             CL("    next();",                    SX_PLAIN),
             CL("  } catch {",                   SX_KEYWORD),
             CL("    res.status(401).json(",      SX_PLAIN),
             CL("      { error: 'Lejárt token' });", SX_STRING),
             CL("  }",                           SX_KEYWORD),
             CL("}",                             SX_KEYWORD),
         ], 10.5),
        (6.7, 1.05, 6.3, 6.2,
         [
             CL("// POST /api/admin/login", SX_COMMENT),
             CL("app.post('/api/admin/login',", SX_FUNC),
             CL("async (req, res) => {",         SX_KEYWORD),
             CL("  const { felhasznalonev,",      SX_PLAIN),
             CL("          jelszo } = req.body;", SX_PLAIN),
             CL("",                              SX_PLAIN),
             CL("  // Admin keresése DB-ben", SX_COMMENT),
             CL("  const admin = db.prepare(",   SX_FUNC),
             CL("    'SELECT * FROM admin",       SX_STRING),
             CL("     WHERE felhasznalonev = ?'", SX_STRING),
             CL("  ).get(felhasznalonev);",       SX_PLAIN),
             CL("",                              SX_PLAIN),
             CL("  if (!admin) return res.status(401)", SX_PLAIN),
             CL("    .json({ error: 'Hibás adatok' });", SX_STRING),
             CL("",                              SX_PLAIN),
             CL("  // Jelszó ellenőrzés bcrypt-tel", SX_COMMENT),
             CL("  const valid = await bcrypt.compare(", SX_FUNC),
             CL("    jelszo, admin.jelszo_hash);",SX_PLAIN),
             CL("  if (!valid) return res.status(401)", SX_PLAIN),
             CL("    .json({ error: 'Hibás adatok' });", SX_STRING),
             CL("",                             SX_PLAIN),
             CL("  // Token visszaküldése", SX_COMMENT),
             CL("  res.json({",                  SX_PLAIN),
             CL("    token: generateToken(admin),", SX_FUNC),
             CL("    admin: { id: admin.id,",    SX_PLAIN),
             CL("      felhasznalonev: admin.felhasznalonev }", SX_PLAIN),
             CL("  });",                         SX_PLAIN),
             CL("});",                           SX_KEYWORD),
         ], 10.5),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# KÓD DIA 10 – Express Server (rendelés létrehozás, DB tranzakció)
# ─────────────────────────────────────────────────────────────────────────────
code_slide(
    "Kód – Express Szerver: Rendelés Létrehozás",
    "POST /api/orders – tranzakció, snapshot árazás, CASCADE kapcsolat  (server.cjs)",
    "server.cjs",
    [
        (0.3, 1.05, 12.7, 6.2,
         [
             CL("// POST /api/orders  –  publikus végpont, nem kell auth", SX_COMMENT),
             CL("app.post('/api/orders', (req, res) => {", SX_FUNC),
             CL("  const { vevo_nev, telefon, email,", SX_PLAIN),
             CL("    iranyitoszam, telepules, utca_hazszam, items } = req.body;", SX_PLAIN),
             CL("",                              SX_PLAIN),
             CL("  // Tranzakció – atomi mentés (rendelés + tételek egyszerre)", SX_COMMENT),
             CL("  const insertOrder = db.transaction(() => {", SX_FUNC),
             CL("    const orderId = generateId();   // crypto.randomBytes(4).toString('hex')", SX_PLAIN),
             CL("",                              SX_PLAIN),
             CL("    // 1. Rendelés fej mentése", SX_COMMENT),
             CL("    db.prepare(`INSERT INTO rendelesek", SX_FUNC),
             CL("      (id, vevo_nev, telefon, email,", SX_STRING),
             CL("       iranyitoszam, telepules, utca_hazszam)", SX_STRING),
             CL("      VALUES (?, ?, ?, ?, ?, ?, ?)`", SX_STRING),
             CL("    ).run(orderId, vevo_nev, telefon, email,", SX_PLAIN),
             CL("           iranyitoszam, telepules, utca_hazszam);", SX_PLAIN),
             CL("",                              SX_PLAIN),
             CL("    // 2. Rendelési tételek mentése (snapshot ár!)", SX_COMMENT),
             CL("    const insertItem = db.prepare(`", SX_FUNC),
             CL("      INSERT INTO rendeles_tetelek", SX_STRING),
             CL("        (id, rendeles_id, termek_nev,", SX_STRING),
             CL("         termek_ar, mennyiseg, tej, cukor)", SX_STRING),
             CL("        VALUES (?, ?, ?, ?, ?, ?, ?)`);", SX_STRING),
             CL("",                              SX_PLAIN),
             CL("    for (const item of items) {", SX_KEYWORD),
             CL("      insertItem.run(generateId(), orderId,", SX_PLAIN),
             CL("        item.termek_nev, item.termek_ar,", SX_PLAIN),
             CL("        item.mennyiseg, item.tej, item.cukor);", SX_PLAIN),
             CL("    }",                          SX_KEYWORD),
             CL("    return orderId;",            SX_PLAIN),
             CL("  });",                          SX_KEYWORD),
             CL("  const newId = insertOrder();  // tranzakció futtatása", SX_FUNC),
             CL("  res.status(201).json({ id: newId, message: 'Rendelés sikeresen létrehozva' });", SX_PLAIN),
             CL("});",                            SX_KEYWORD),
         ], 11),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# KÓD DIA 11 – SQLite séma (CREATE TABLE)
# ─────────────────────────────────────────────────────────────────────────────
code_slide(
    "Kód – Adatbázis Séma",
    "SQLite CREATE TABLE utasítások, indexek, CASCADE törlés  (schema.sql)",
    "schema.sql",
    [
        (0.3, 1.05, 6.2, 6.2,
         [
             CL("-- Termék katalógus", SX_COMMENT),
             CL("CREATE TABLE IF NOT EXISTS termekek (", SX_SQL_KW),
             CL("    id      TEXT    PRIMARY KEY,", SX_PLAIN),
             CL("    nev     TEXT    NOT NULL,",    SX_PLAIN),
             CL("    ar      INTEGER NOT NULL,",    SX_PLAIN),
             CL("    kep_url TEXT",                 SX_PLAIN),
             CL(");",                               SX_SQL_KW),
             CL("",                                SX_PLAIN),
             CL("-- Kezdő termékek", SX_COMMENT),
             CL("INSERT OR IGNORE INTO termekek", SX_SQL_KW),
             CL("  (id, nev, ar, kep_url) VALUES", SX_PLAIN),
             CL("('a1b2c3d4','Cappuccino', 850,",  SX_STRING),
             CL("  '/images/cappuccino.jpg'),",     SX_STRING),
             CL("('b2c3d4e5','Espresso',   650,",  SX_STRING),
             CL("  '/images/espresso.webp'),",      SX_STRING),
             CL("('c3d4e5f6','Ristretto',  850,",  SX_STRING),
             CL("  '/images/ristretto.jpg'),",      SX_STRING),
             CL("('d4e5f6a7','Latte',      900,",  SX_STRING),
             CL("  '/images/latte.jpg'),",          SX_STRING),
             CL("('e5f6a7b8','Americano',  700,",  SX_STRING),
             CL("  '/images/americano.jpg'),",      SX_STRING),
             CL("('f6a7b8c9','Doppio',     800,",  SX_STRING),
             CL("  '/images/doppio.webp');",        SX_STRING),
         ], 10.5),
        (6.7, 1.05, 6.3, 6.2,
         [
             CL("-- Rendelési tételek (FK + CASCADE)", SX_COMMENT),
             CL("CREATE TABLE IF NOT EXISTS rendeles_tetelek (", SX_SQL_KW),
             CL("    id          TEXT    PRIMARY KEY,", SX_PLAIN),
             CL("    rendeles_id TEXT    NOT NULL,",    SX_PLAIN),
             CL("    termek_nev  TEXT    NOT NULL,",    SX_PLAIN),
             CL("    termek_ar   INTEGER NOT NULL,",    SX_PLAIN),
             CL("    mennyiseg   INTEGER NOT NULL,",    SX_PLAIN),
             CL("    tej         TEXT    NOT NULL,",    SX_PLAIN),
             CL("    cukor       TEXT    NOT NULL,",    SX_PLAIN),
             CL("    FOREIGN KEY(rendeles_id)",         SX_SQL_KW),
             CL("      REFERENCES rendelesek(id)",      SX_PLAIN),
             CL("      ON DELETE CASCADE",              SX_SQL_KW),
             CL(");",                                   SX_SQL_KW),
             CL("",                                    SX_PLAIN),
             CL("-- Teljesítmény indexek", SX_COMMENT),
             CL("CREATE INDEX IF NOT EXISTS",           SX_SQL_KW),
             CL("  idx_rendelesek_megrendelve",         SX_PLAIN),
             CL("  ON rendelesek(megrendelve);",         SX_PLAIN),
             CL("CREATE INDEX IF NOT EXISTS",           SX_SQL_KW),
             CL("  idx_rendelesek_postazva",            SX_PLAIN),
             CL("  ON rendelesek(postazva);",           SX_PLAIN),
             CL("CREATE INDEX IF NOT EXISTS",           SX_SQL_KW),
             CL("  idx_rendeles_tetelek_rendeles_id",   SX_PLAIN),
             CL("  ON rendeles_tetelek(rendeles_id);",  SX_PLAIN),
         ], 10.5),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# KÓD DIA 12 – Tesztek (cart.test.ts + api.test.ts)
# ─────────────────────────────────────────────────────────────────────────────
code_slide(
    "Kód – Vitest Tesztek",
    "Kosár store tesztek és API mock tesztek  (cart.test.ts / api.test.ts)",
    "tests/cart.test.ts + api.test.ts",
    [
        (0.3, 1.05, 6.2, 6.2,
         [
             CL("// cart.test.ts – Kosár tesztek", SX_COMMENT),
             CL("import { describe, it, expect,", SX_KEYWORD),
             CL("         beforeEach } from 'vitest';", SX_KEYWORD),
             CL("import { get } from 'svelte/store';", SX_KEYWORD),
             CL("import { cartItems, cartTotal,",  SX_KEYWORD),
             CL("         addToCart } from '../src/lib/cart';", SX_KEYWORD),
             CL("",                               SX_PLAIN),
             CL("describe('Cart Store', () => {", SX_FUNC),
             CL("  beforeEach(() => cartItems.set([]));", SX_PLAIN),
             CL("",                               SX_PLAIN),
             CL("  it('kezdetben üres a kosár', () => {", SX_FUNC),
             CL("    expect(get(cartItems)).toEqual([]);", SX_PLAIN),
             CL("  });",                           SX_KEYWORD),
             CL("",                               SX_PLAIN),
             CL("  it('addToCart hozzáad egy terméket', () => {", SX_FUNC),
             CL("    addToCart(mockProduct, 2,",   SX_PLAIN),
             CL("      { milk: 'Zabtej', sugar: '1 kanál' });", SX_STRING),
             CL("    const items = get(cartItems);", SX_PLAIN),
             CL("    expect(items).toHaveLength(1);", SX_PLAIN),
             CL("    expect(items[0].quantity).toBe(2);", SX_PLAIN),
             CL("  });",                           SX_KEYWORD),
             CL("",                               SX_PLAIN),
             CL("  it('cartTotal helyes összeget számol', () => {", SX_FUNC),
             CL("    addToCart(mockProduct, 2, mockModifiers);", SX_PLAIN),
             CL("    // 850 * 2 = 1700",           SX_COMMENT),
             CL("    expect(get(cartTotal)).toBe(1700);", SX_PLAIN),
             CL("  });",                           SX_KEYWORD),
             CL("});",                             SX_KEYWORD),
         ], 10.5),
        (6.7, 1.05, 6.3, 6.2,
         [
             CL("// api.test.ts – API mock tesztek", SX_COMMENT),
             CL("import { vi } from 'vitest';",    SX_KEYWORD),
             CL("",                               SX_PLAIN),
             CL("// fetch és localStorage mock",   SX_COMMENT),
             CL("const mockFetch = vi.fn();",      SX_FUNC),
             CL("vi.stubGlobal('fetch', mockFetch);", SX_PLAIN),
             CL("vi.stubGlobal('localStorage', {", SX_PLAIN),
             CL("  getItem:    (k) => store[k] ?? null,", SX_PLAIN),
             CL("  setItem:    (k, v) => store[k] = v,", SX_PLAIN),
             CL("  removeItem: (k) => delete store[k]", SX_PLAIN),
             CL("});",                             SX_KEYWORD),
             CL("",                               SX_PLAIN),
             CL("describe('Token Management', () => {", SX_FUNC),
             CL("  it('setToken elmenti a tokent', () => {", SX_FUNC),
             CL("    setToken('test-token-123');", SX_PLAIN),
             CL("    expect(store['admin_token'])", SX_PLAIN),
             CL("      .toBe('test-token-123');",  SX_STRING),
             CL("  });",                           SX_KEYWORD),
             CL("",                               SX_PLAIN),
             CL("  it('clearToken törli a tokent', () => {", SX_FUNC),
             CL("    store['admin_token'] = 'xyz';", SX_PLAIN),
             CL("    clearToken();",               SX_PLAIN),
             CL("    expect(store['admin_token'])", SX_PLAIN),
             CL("      .toBeUndefined();",         SX_PLAIN),
             CL("  });",                           SX_KEYWORD),
             CL("});",                             SX_KEYWORD),
         ], 10.5),
    ]
)


# ─────────────────────────────────────────────────────────────────────────────
# MENTÉS
# ─────────────────────────────────────────────────────────────────────────────
output_path = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "NeoCoffee_Webáruház_Prezentáció.pptx"
)
prs.save(output_path)
print(f"✅  Prezentáció mentve: {output_path}")
print(f"    Diák száma: {len(prs.slides)}")
